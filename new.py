# Dictionary to hold the topics and sub-topics
data = {
    "Data Science": {
        "Sub Topics": ["Data Mining", "Big Data"],
        "Intro": {
            "Data Mining": "Data mining is the process of extracting valuable information from large datasets.",
            "Big Data": "Big data refers to extremely large datasets that can be analyzed to reveal patterns, trends, and associations."
        }
    },
    "A.I": {
        "Sub Topics": ["Machine Learning", "Computer Vision"],
        "Intro": {
            "Machine Learning": "Machine learning is a method of teaching computers to learn from data without being explicitly programmed.",
            "Computer Vision": "Computer vision is a field of artificial intelligence that trains computers to interpret and understand visual data from the world around them."
        }
    },
    "ML": {
        "Sub Topics": ["Deep Learning", "Neural Networks"],
        "Intro": {
            "Deep Learning": "Deep learning is a type of machine learning that uses artificial neural networks to model and solve complex problems.",
            "Neural Networks": "Neural networks are a type of machine learning algorithm that are modeled after the structure and function of the human brain."
        }
    },
    "CSE": {
        "Sub Topics": ["Operating Systems", "Computer Networks"],
        "Intro": {
            "Operating Systems": "An operating system is a software that manages computer hardware and software resources and provides common services for computer programs.",
            "Computer Networks": "A computer network is a group of computers and other devices connected together to share resources and communicate with each other."
        }
    },
    "IOT": {
        "Sub Topics": ["Wireless Sensor Networks", "Smart Home"],
        "Intro": {
            "Wireless Sensor Networks": "Wireless sensor networks are networks of spatially distributed sensors that communicate with each other and send data to a central location for analysis.",
            "Smart Home": "A smart home is a residence that uses internet-connected devices to control and automate lighting, heating, and other appliances."
        }
    }
}

# Get user input for topic and sub-topic
topic_choice = input("Select a topic from the following: Data Science, A.I, ML, CSE, IOT: ")
sub_topic_choice = input(f"Select a sub-topic from the following {data[topic_choice]['Sub Topics']}: ")

# Retrieve introduction of the selected sub-topic
intro = data[topic_choice]['Intro'][sub_topic_choice]

# Print the introduction


import collections.abc 
from pptx import Presentation
from pptx.util import Inches



# Create presentation
prs = Presentation()

 # define slidelayouts 
slide = prs.slides.add_slide(prs.slide_layouts[0])


# title slide
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = sub_topic_choice
subtitle.text = " "

# Introduction page
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

# add a title to the slide
title_shape = slide.shapes.title
title_shape.text = "Introduction"

# add some content to the slide
body_shape = slide.shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = intro


# end slide
slide3 = prs.slides.add_slide(prs.slide_layouts[0])
slide = slide3.shapes.title
subtitle = slide3.placeholders[1]
slide.text = "End of Presentation"
subtitle.text = "Thank You"

# export and save the presentation
prs.save(sub_topic_choice +'.pptx')

print("Your ppt is Created...")